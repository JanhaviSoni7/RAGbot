"""ragbot/ingestion/loader.py

Loads PDF, plain-text, markdown, office documents, and image documents.

Returns a list of `RawDocument` objects, each carrying:
  - source      : original file path
  - page_no     : 0-indexed page number (0 for single-page sources)
  - content     : extracted text
  - metadata    : dict with doc-level / page-level metadata
"""

from __future__ import annotations

import io
import csv
import json
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional

from config import SUPPORTED_EXTENSIONS
from ragbot.utils.logger import get_logger

log = get_logger(__name__)

# ── Optional heavy imports ────────────────────────────────────────────────────
try:
    import fitz  # PyMuPDF
    import pymupdf4llm
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False
    log.warning("PyMuPDF not installed — PDF support disabled.")

try:
    from PIL import Image
    import pytesseract
    HAS_OCR = True
except ImportError:
    HAS_OCR = False
    log.warning("pytesseract / Pillow not installed — OCR support disabled.")

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    log.warning("python-docx not installed — DOCX support disabled.")

try:
    from openpyxl import load_workbook
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False
    log.warning("openpyxl not installed — XLSX support disabled.")

try:
    import xlrd
    HAS_XLRD = True
except ImportError:
    HAS_XLRD = False
    log.warning("xlrd not installed — XLS support disabled.")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    log.warning("python-pptx not installed — PPTX support disabled.")


@dataclass
class RawDocument:
    source:   str
    page_no:  int
    content:  str
    metadata: dict = field(default_factory=dict)


# ─────────────────────────────────────────────────────────────────────────────
# Internal helpers
# ─────────────────────────────────────────────────────────────────────────────

def _load_pdf(path: Path) -> List[RawDocument]:
    """Extract text/markdown from each PDF page using PyMuPDF4LLM."""
    if not HAS_PYMUPDF:
        raise RuntimeError("PyMuPDF is required to load PDF files.")

    docs: List[RawDocument] = []
    try:
        # pymupdf4llm gives us clean markdown per page
        md_pages = pymupdf4llm.to_markdown(str(path), page_chunks=True)
        for page_info in md_pages:
            text = page_info.get("text", "").strip()
            if not text:
                # Fall back to raw fitz text extraction
                pdf = fitz.open(str(path))
                pg  = page_info.get("page", 0)
                text = pdf[pg].get_text("text").strip()
                pdf.close()
            if text:
                docs.append(RawDocument(
                    source   = str(path),
                    page_no  = page_info.get("page", 0),
                    content  = text,
                    metadata = {
                        "file_name": path.name,
                        "file_type": "pdf",
                        "page":      page_info.get("page", 0),
                        "total_pages": page_info.get("total_pages", 1),
                    },
                ))
    except Exception as exc:
        log.error("PDF load failed: %s — %s", path, exc)

    # If pymupdf4llm gave nothing, fall back page-by-page
    if not docs and HAS_PYMUPDF:
        try:
            pdf = fitz.open(str(path))
            for pg_no, page in enumerate(pdf):
                text = page.get_text("text").strip()
                if text:
                    docs.append(RawDocument(
                        source   = str(path),
                        page_no  = pg_no,
                        content  = text,
                        metadata = {"file_name": path.name, "file_type": "pdf",
                                    "page": pg_no, "total_pages": len(pdf)},
                    ))
                else:
                    # Page might be image-based — try OCR
                    if HAS_OCR:
                        pix  = page.get_pixmap(dpi=300)
                        img  = Image.open(io.BytesIO(pix.tobytes("png")))
                        text = pytesseract.image_to_string(img).strip()
                        if text:
                            docs.append(RawDocument(
                                source   = str(path),
                                page_no  = pg_no,
                                content  = text,
                                metadata = {"file_name": path.name, "file_type": "pdf_ocr",
                                            "page": pg_no, "total_pages": len(pdf)},
                            ))
            pdf.close()
        except Exception as exc:
            log.error("Fallback PDF load failed: %s — %s", path, exc)

    log.info("Loaded PDF '%s': %d page(s)", path.name, len(docs))
    return docs


def _load_text(path: Path) -> List[RawDocument]:
    """Load plain text / markdown files."""
    try:
        text = path.read_text(encoding="utf-8", errors="ignore").strip()
    except Exception as exc:
        log.error("Text load failed: %s — %s", path, exc)
        return []

    if not text:
        return []

    return [RawDocument(
        source   = str(path),
        page_no  = 0,
        content  = text,
        metadata = {"file_name": path.name, "file_type": path.suffix.lstrip(".")},
    )]


def _load_image(path: Path) -> List[RawDocument]:
    """OCR-extract text from image documents."""
    if not HAS_OCR:
        raise RuntimeError("pytesseract + Pillow required for image loading.")

    try:
        img  = Image.open(str(path))
        text = pytesseract.image_to_string(img).strip()
    except Exception as exc:
        log.error("Image OCR failed: %s — %s", path, exc)
        return []

    if not text:
        log.warning("OCR extracted no text from '%s'.", path.name)
        return []

    return [RawDocument(
        source   = str(path),
        page_no  = 0,
        content  = text,
        metadata = {"file_name": path.name, "file_type": "image_ocr"},
    )]


def _load_docx(path: Path) -> List[RawDocument]:
    """Extract text from DOCX files."""
    if not HAS_DOCX:
        raise RuntimeError("python-docx is required to load DOCX files.")

    try:
        doc = DocxDocument(str(path))
        text = "\n".join(para.text for para in doc.paragraphs).strip()
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += "\n" + cell.text
        
        text = text.strip()
    except Exception as exc:
        log.error("DOCX load failed: %s — %s", path, exc)
        return []

    if not text:
        log.warning("No text extracted from DOCX '%s'.", path.name)
        return []

    return [RawDocument(
        source   = str(path),
        page_no  = 0,
        content  = text,
        metadata = {"file_name": path.name, "file_type": "docx"},
    )]


def _load_xlsx(path: Path) -> List[RawDocument]:
    """Extract text from XLSX files."""
    if not HAS_OPENPYXL:
        raise RuntimeError("openpyxl is required to load XLSX files.")

    try:
        workbook = load_workbook(str(path), data_only=True)
        text_parts = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"Sheet: {sheet_name}")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(
                    str(cell) if cell is not None else ""
                    for cell in row
                ).strip()
                if row_text:
                    text_parts.append(row_text)
        
        text = "\n".join(text_parts).strip()
    except Exception as exc:
        log.error("XLSX load failed: %s — %s", path, exc)
        return []

    if not text:
        log.warning("No data extracted from XLSX '%s'.", path.name)
        return []

    return [RawDocument(
        source   = str(path),
        page_no  = 0,
        content  = text,
        metadata = {"file_name": path.name, "file_type": "xlsx"},
    )]


def _load_xls(path: Path) -> List[RawDocument]:
    """Extract text from XLS files."""
    if not HAS_XLRD:
        raise RuntimeError("xlrd is required to load XLS files.")

    try:
        workbook = xlrd.open_workbook(str(path), data_only=True)
        text_parts = []
        
        for sheet_idx in range(workbook.nsheets):
            sheet = workbook.sheet_by_index(sheet_idx)
            text_parts.append(f"Sheet: {sheet.name}")
            
            for row_idx in range(sheet.nrows):
                row_values = sheet.row_values(row_idx)
                row_text = " | ".join(
                    str(cell) if cell else ""
                    for cell in row_values
                ).strip()
                if row_text:
                    text_parts.append(row_text)
        
        text = "\n".join(text_parts).strip()
    except Exception as exc:
        log.error("XLS load failed: %s — %s", path, exc)
        return []

    if not text:
        log.warning("No data extracted from XLS '%s'.", path.name)
        return []

    return [RawDocument(
        source   = str(path),
        page_no  = 0,
        content  = text,
        metadata = {"file_name": path.name, "file_type": "xls"},
    )]


def _load_pptx(path: Path) -> List[RawDocument]:
    """Extract text from PPTX files."""
    if not HAS_PPTX:
        raise RuntimeError("python-pptx is required to load PPTX files.")

    try:
        presentation = Presentation(str(path))
        text_parts = []
        
        for slide_idx, slide in enumerate(presentation.slides):
            text_parts.append(f"Slide {slide_idx + 1}:")
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        text_parts.append(shape.text)
        
        text = "\n".join(text_parts).strip()
    except Exception as exc:
        log.error("PPTX load failed: %s — %s", path, exc)
        return []

    if not text:
        log.warning("No text extracted from PPTX '%s'.", path.name)
        return []

    return [RawDocument(
        source   = str(path),
        page_no  = 0,
        content  = text,
        metadata = {"file_name": path.name, "file_type": "pptx"},
    )]


def _load_csv(path: Path) -> List[RawDocument]:
    """Extract text from CSV files."""
    try:
        with open(str(path), 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            rows = list(reader)
        
        if not rows:
            log.warning("CSV file '%s' is empty.", path.name)
            return []
        
        text_parts = []
        for idx, row in enumerate(rows):
            row_text = " | ".join(str(cell).strip() for cell in row).strip()
            if row_text:
                text_parts.append(row_text)
        
        text = "\n".join(text_parts).strip()
    except Exception as exc:
        log.error("CSV load failed: %s — %s", path, exc)
        return []

    if not text:
        return []

    return [RawDocument(
        source   = str(path),
        page_no  = 0,
        content  = text,
        metadata = {"file_name": path.name, "file_type": "csv", "rows": len(rows)},
    )]


def _load_json(path: Path) -> List[RawDocument]:
    """Extract text from JSON files."""
    try:
        with open(str(path), 'r', encoding='utf-8', errors='ignore') as f:
            data = json.load(f)
        
        # Convert JSON structure to readable text
        text = _json_to_text(data)
        text = text.strip()
    except Exception as exc:
        log.error("JSON load failed: %s — %s", path, exc)
        return []

    if not text:
        log.warning("No content extracted from JSON '%s'.", path.name)
        return []

    return [RawDocument(
        source   = str(path),
        page_no  = 0,
        content  = text,
        metadata = {"file_name": path.name, "file_type": "json"},
    )]


def _json_to_text(data: any, indent: int = 0) -> str:
    """Recursively convert JSON data structure to readable text."""
    text_parts = []
    prefix = "  " * indent
    
    if isinstance(data, dict):
        for key, value in data.items():
            if isinstance(value, (dict, list)):
                text_parts.append(f"{prefix}{key}:")
                text_parts.append(_json_to_text(value, indent + 1))
            else:
                text_parts.append(f"{prefix}{key}: {value}")
    elif isinstance(data, list):
        for idx, item in enumerate(data):
            if isinstance(item, (dict, list)):
                text_parts.append(f"{prefix}[{idx}]:")
                text_parts.append(_json_to_text(item, indent + 1))
            else:
                text_parts.append(f"{prefix}[{idx}]: {item}")
    else:
        text_parts.append(f"{prefix}{data}")
    
    return "\n".join(text_parts)


# ─────────────────────────────────────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────────────────────────────────────

def load_document(path: str | Path) -> List[RawDocument]:
    """Load a single document and return its pages as RawDocuments."""
    path = Path(path)
    ext  = path.suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: '{ext}'. "
                         f"Supported: {SUPPORTED_EXTENSIONS}")

    if ext == ".pdf":
        return _load_pdf(path)
    elif ext in {".png", ".jpg", ".jpeg", ".tiff"}:
        return _load_image(path)
    elif ext == ".docx":
        return _load_docx(path)
    elif ext == ".xlsx":
        return _load_xlsx(path)
    elif ext == ".xls":
        return _load_xls(path)
    elif ext == ".pptx":
        return _load_pptx(path)
    elif ext == ".csv":
        return _load_csv(path)
    elif ext == ".json":
        return _load_json(path)
    else:  # .txt, .md
        return _load_text(path)


def load_documents(paths: List[str | Path]) -> List[RawDocument]:
    """Load multiple documents; skip problematic files with a warning."""
    results: List[RawDocument] = []
    for p in paths:
        try:
            docs = load_document(p)
            results.extend(docs)
            log.info("Loaded '%s' → %d unit(s)", Path(p).name, len(docs))
        except Exception as exc:
            log.error("Skipping '%s': %s", p, exc)
    return results
